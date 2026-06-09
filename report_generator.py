import io
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from branding import COLORS, get_rgb
from data_parser import STATS
from insights import generate_insights

def safe_float(val):
    try:
        return float(val)
    except (ValueError, TypeError):
        return 0.0

def create_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

# ── MATPLOTLIB STYLE ENGINE (Union 2 Inspired) ─────────────────
def generate_chart_stream(values, labels, colors, title):
    """Generates an explicit bar chart directly inside an in-memory byte stream."""
    fig, ax = plt.subplots(figsize=(5.5, 4.0), dpi=200)
    fig.patch.set_facecolor('#' + COLORS["WHITE"])
    ax.set_facecolor('#' + COLORS["WHITE"])
    
    bars = ax.bar(labels, values, color=colors, width=0.5, edgecolor=None)
    ax.set_title(title.upper(), fontsize=11, fontweight='bold', pad=15, color='#' + COLORS["BLACK"])
    
    # Clean, minimalist presentation axes
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_color('#' + COLORS["SILVER"])
    ax.yaxis.grid(True, linestyle='--', alpha=0.5, color='#' + COLORS["SILVER"])
    ax.set_axisbelow(True)
    ax.tick_params(axis='x', colors='#' + COLORS["BLACK"], labelsize=9)
    ax.tick_params(axis='y', colors='#' + COLORS["SILVER"], labelsize=8)

    # Inline text labels placed above bars to echo professional scouting reports
    for bar in bars:
        height = bar.get_height()
        ax.annotate(f'{height:.2f}',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 4),  
                    textcoords="offset points",
                    ha='center', va='bottom', fontsize=10, fontweight='bold', color='#' + COLORS["BLACK"])

    plt.tight_layout()
    img_buf = io.BytesIO()
    plt.savefig(img_buf, format='png', facecolor=fig.get_facecolor(), edgecolor='none')
    img_buf.seek(0)
    plt.close(fig)
    return img_buf

# ── SLIDE BUILDERS ─────────────────────────────────────────────
def add_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # High-contrast premium dark canvas
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["BLACK"])

    # Golden Header Stripe Accent Accent
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = get_rgb(COLORS["GOLD"])
    stripe.line.fill.background()

    # Club Heading
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BROOKLYN FC"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = get_rgb(COLORS["GOLD"])
    p.font.name = "Arial Black"
    p.alignment = PP_ALIGN.CENTER

    # Sub-heading
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.33), Inches(0.6))
    p2 = sub_box.text_frame.paragraphs[0]
    p2.text = "MATCH ANALYSIS REPORT"
    p2.font.size = Pt(16)
    p2.font.color.rgb = get_rgb(COLORS["WHITE"])
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER

    # Score Badge
    score_box = slide.shapes.add_textbox(Inches(5.16), Inches(4.0), Inches(3.0), Inches(0.7))
    tf_score = score_box.text_frame
    p_score = tf_score.paragraphs[0]
    p_score.text = f" {data['score']} "
    p_score.font.size = Pt(32)
    p_score.font.bold = True
    p_score.font.color.rgb = get_rgb(COLORS["BLACK"])
    p_score.font.name = "Arial Black"
    p_score.alignment = PP_ALIGN.CENTER
    
    # Custom score block shape background styling
    score_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.16), Inches(4.0), Inches(3.0), Inches(0.7))
    score_bg.fill.solid()
    score_bg.fill.fore_color.rgb = get_rgb(COLORS["GOLD"])
    score_bg.line.fill.background()
    # Move shape behind text box layers
    slide.shapes._spTree.remove(score_bg._element)
    slide.shapes._spTree.insert(2, score_bg._element)

    # Context Metadata Frame
    meta_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.33), Inches(1.0))
    tf_meta = meta_box.text_frame
    
    p3 = tf_meta.paragraphs[0]
    p3.text = f"BKFC vs {data['opponent_name'].upper()}"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = get_rgb(COLORS["WHITE"])
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf_meta.add_paragraph()
    p4.text = f"{data['match_date']}  |  {data['competition']}"
    p4.font.size = Pt(13)
    p4.font.color.rgb = get_rgb(COLORS["SILVER"])
    p4.alignment = PP_ALIGN.CENTER

def add_summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    # Header banner block setup
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = get_rgb(COLORS["BLACK"])
    banner.line.fill.background()

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.9), Inches(13.33), Inches(0.05))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = get_rgb(COLORS["GOLD"])
    stripe.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.33), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "MATCH PERFORMANCE OVERVIEW"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = get_rgb(COLORS["GOLD"])
    p.font.name = "Arial Black"

    # Core Metric Matrix Table Layout Engine
    table_shape = slide.shapes.add_table(rows=7, cols=3, left=Inches(1.66), top=Inches(1.8), width=Inches(10), height=Inches(4.5))
    table = table_shape.table

    headers = ["KEY METRIC PERFORMANCE", "BROOKLYN FC", data['opponent_name'].upper()]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = get_rgb(COLORS["BLACK"])
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = get_rgb(COLORS["GOLD"])
        p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER

    summary_metrics = [
        ("Goals", 6),
        ("Expected Goals (xG)", 7),
        ("Shots Taken", 8),
        ("Shots on Target", 9),
        ("Possession %", 14),
        ("Pass Accuracy %", 13)
    ]

    for row_idx, (label, col) in enumerate(summary_metrics, start=1):
        table.cell(row_idx, 0).text = label
        table.cell(row_idx, 1).text = f"{safe_float(data['match_bkfc'][col]):.2f}"
        table.cell(row_idx, 2).text = f"{safe_float(data['match_opp'][col]):.2f}"
        
        for col_idx in range(3):
            cell = table.cell(row_idx, col_idx)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = get_rgb(COLORS["BLACK"])
            if col_idx > 0:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True

def add_insights_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    # Uniform header structures matching layout controls
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = get_rgb(COLORS["BLACK"])
    banner.line.fill.background()

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.9), Inches(13.33), Inches(0.05))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = get_rgb(COLORS["GOLD"])
    stripe.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.33), Inches(0.6))
    title_box.text_frame.paragraphs[0].text = "STRATEGIC MATCH INSIGHTS"
    title_box.text_frame.paragraphs[0].font.size = Pt(22)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = get_rgb(COLORS["GOLD"])
    title_box.text_frame.paragraphs[0].font.name = "Arial Black"

    # Parsing telemetry arrays for performance evaluations
    stats_list = []
    for s in STATS:
        stats_list.append({
            "label": s["label"],
            "match": safe_float(data["match_bkfc"][s["col"]]),
            "season": safe_float(data["bkfc_season_avg"][s["col"]])
        })

    insights = generate_insights(stats_list)

    if not insights:
        insights = ["Performance patterns tracked closely within expected structural targets across all phases."]

    y_pos = Inches(1.8)
    for insight_text in insights:
        # Drawing structured container cards to isolate discrete bullet layers
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y_pos, Inches(10.33), Inches(0.65))
        card.fill.solid()
        card.fill.fore_color.rgb = get_rgb(COLORS["DARK_GRAY"])
        card.line.color.rgb = get_rgb(COLORS["GOLD"])
        card.line.width = Pt(1)

        box = slide.shapes.add_textbox(Inches(1.7), y_pos + Inches(0.05), Inches(10), Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"▪  {insight_text}"
        p.font.size = Pt(13)
        p.font.color.rgb = get_rgb(COLORS["WHITE"])
        p.font.bold = True
        
        y_pos += Inches(0.85)

def add_stat_slide(prs, data, label, col):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = get_rgb(COLORS["WHITE"])

    # Section Headers
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.7))
    banner.fill.solid()
    banner.fill.fore_color.rgb = get_rgb(COLORS["BLACK"])
    banner.line.fill.background()

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.7), Inches(13.33), Inches(0.04))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = get_rgb(COLORS["GOLD"])
    stripe.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = get_rgb(COLORS["GOLD"])
    p.font.name = "Arial Black"

    # Context Section Subheaders
    sub_match = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5.8), Inches(0.3))
    p_sm = sub_match.text_frame.paragraphs[0]
    p_sm.text = "THIS MATCH ANALYSIS"
    p_sm.font.size = Pt(11)
    p_sm.font.bold = True
    p_sm.font.color.rgb = get_rgb(COLORS["SILVER"])
    p_sm.alignment = PP_ALIGN.CENTER

    sub_season = slide.shapes.add_textbox(Inches(7.0), Inches(0.9), Inches(5.8), Inches(0.3))
    p_ss = sub_season.text_frame.paragraphs[0]
    p_ss.text = "SEASON BASELINE BENCHMARKS"
    p_ss.font.size = Pt(11)
    p_ss.font.bold = True
    p_ss.font.color.rgb = get_rgb(COLORS["SILVER"])
    p_ss.alignment = PP_ALIGN.CENTER

    # Structural Grid Segment Partition Accent Line
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.65), Inches(1.3), Inches(0.02), Inches(5.2))
    divider.fill.solid()
    divider.fill.fore_color.rgb = get_rgb(COLORS["GRID"])
    divider.line.fill.background()

    # Color Mapping Rules (BKFC Gold, League/Opps Silvers/Grays)
    match_colors = ['#' + COLORS["GOLD"], '#' + COLORS["DARK_GRAY"]]
    season_colors = ['#' + COLORS["GOLD"], '#' + COLORS["SILVER"], '#' + COLORS["BLACK"]]

    # Chart 1: This Match View (BKFC vs Match Opponent)
    match_vals = [safe_float(data["match_bkfc"][col]), safe_float(data["match_opp"][col])]
    match_labels = ["BKFC", data["opponent_name"]]
    match_chart_stream = generate_chart_stream(match_vals, match_labels, match_colors, "Match Comparison")

    # Chart 2: Season Performance Vector (BKFC Avg, League Avg, Opponent Baseline)
    opp_col_offset = col - 1
    opp_season_val = safe_float(data["opp_season_avg"][opp_col_offset]) if opp_col_offset in data["opp_season_avg"] else 0.0
    
    season_vals = [
        safe_float(data["bkfc_season_avg"][col]),
        safe_float(data["all_opp_avg"][col]),
        opp_season_val
    ]
    season_labels = ["BKFC Baseline", "Opponent Average", f"{data['opponent_name']} Baseline"]
    season_chart_stream = generate_chart_stream(season_vals, season_labels, season_colors, "Historical Performance Context")

    # Injecting Data Streams onto Presentation Canvas Maps
    slide.shapes.add_picture(match_chart_stream, Inches(0.6), Inches(1.5), Inches(5.5), Inches(4.8))
    slide.shapes.add_picture(season_chart_stream, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.8))

    # Clean systemic slide margins footer branding
    footer = slide.shapes.add_textbox(Inches(0), Inches(7.1), Inches(13.33), Inches(0.3))
    pf = footer.text_frame.paragraphs[0]
    pf.text = "Brooklyn FC Technical Scouting Intelligence"
    pf.font.size = Pt(8)
    pf.font.italic = True
    pf.font.color.rgb = get_rgb(COLORS["SILVER"])
    pf.alignment = PP_ALIGN.CENTER

# ── MAIN PIPELINE ENTRY EXPORT ENGINE ────────────────────────
def generate_report(data):
    """Compiles all individual presentation elements into an in-memory BytesIO track."""
    prs = create_prs()

    add_title_slide(prs, data)
    add_summary_slide(prs, data)
    add_insights_slide(prs, data)

    for stat in STATS:
        add_stat_slide(prs, data, stat["label"], stat["col"])

    prs_stream = io.BytesIO()
    prs.save(prs_stream)
    prs_stream.seek(0)
    return prs_stream